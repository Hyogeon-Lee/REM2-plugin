"""PPTX → Markdown 변환 어댑터.

전략:
- 슬라이드별 ## 헤더 (슬라이드 N)
- shape를 (top, left) 기준으로 정렬 → 위→아래, 좌→우 시각적 순서 보존
- 텍스트 박스: 텍스트 추출 (들여쓰기 레벨로 리스트 처리)
- Picture shape: 이미지 추출 + 해당 위치에 ![]() 삽입
- 표 shape: 마크다운 표
- 발표자 노트: 슬라이드 끝에 "**Notes:**" 블록으로 추가
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..common import ImageSaver, make_image_markdown


def convert(input_path: Path, output_md_path: Path) -> dict:
    prs = Presentation(str(input_path))
    saver = ImageSaver(output_md_path, output_md_path.stem, image_dir_name="images")
    md_lines: list[str] = []
    warnings: list[str] = []

    md_lines.append(f"# {output_md_path.stem}")

    for idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"\n## 슬라이드 {idx}")

        # shape를 시각적 순서로 정렬 (top → left)
        shapes = list(slide.shapes)
        shapes.sort(
            key=lambda s: (
                s.top if s.top is not None else 0,
                s.left if s.left is not None else 0,
            )
        )

        for shape in shapes:
            block = _render_shape(shape, saver, warnings)
            if block:
                md_lines.append(block)

        # 발표자 노트
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                md_lines.append("\n**Notes:**\n")
                md_lines.append(notes_text)

    output_md_path.write_text("\n\n".join(md_lines) + "\n", encoding="utf-8")
    warnings.extend(saver.warnings)
    return {"image_count": saver.saved_count, "warnings": warnings}


def _render_shape(shape, saver: ImageSaver, warnings: list[str]) -> str:
    """shape 하나를 마크다운 블록으로 변환."""
    try:
        st = shape.shape_type
    except Exception:
        st = None

    # 그룹: 재귀
    if st == MSO_SHAPE_TYPE.GROUP:
        parts = []
        for child in shape.shapes:
            block = _render_shape(child, saver, warnings)
            if block:
                parts.append(block)
        return "\n\n".join(parts)

    # 이미지
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            blob = image.blob
            name = getattr(image, "filename", None) or f"slide_image.{image.ext}"
            ext = image.ext.lstrip(".")
            # PPTX shape 크기 (EMU → mm). 914400 EMU = 1 inch = 25.4 mm → mm = emu / 36000
            size_mm = None
            if shape.width and shape.height:
                size_mm = (shape.width / 36000.0, shape.height / 36000.0)
            rel_path = saver.save(
                blob, suggested_name=name, ext=ext, physical_size_mm=size_mm
            )
            return make_image_markdown(rel_path, alt="image")
        except Exception as e:
            warnings.append(f"이미지 추출 실패: {e}")
            return ""

    # 표
    if shape.has_table:
        return _render_table(shape.table)

    # 텍스트 프레임
    if shape.has_text_frame:
        return _render_text_frame(shape.text_frame)

    return ""


def _render_text_frame(tf) -> str:
    lines = []
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            # runs가 비어있을 때 paragraph.text fallback
            text = (para.text or "").strip()
            if not text:
                continue
        # 들여쓰기 레벨 → 리스트
        level = para.level or 0
        if level > 0 or _looks_like_bullet(para):
            indent = "  " * level
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def _looks_like_bullet(para) -> bool:
    """문단이 불릿/번호 매기기인지 휴리스틱 판단."""
    try:
        pPr = para._pPr
        if pPr is None:
            return False
        # buChar (불릿 문자) 또는 buAutoNum (자동 번호)가 있으면 리스트
        return any(
            etree_tag(child) in ("buChar", "buAutoNum") for child in pPr
        )
    except Exception:
        return False


def etree_tag(el) -> str:
    """네임스페이스를 제거한 태그 이름."""
    tag = el.tag
    if "}" in tag:
        return tag.split("}", 1)[1]
    return tag


def _render_table(table) -> str:
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            text = (cell.text or "").replace("\n", "<br>").replace("|", r"\|").strip()
            cells.append(text or " ")
        rows.append(cells)
    if not rows:
        return ""
    n_cols = max(len(r) for r in rows)
    rows = [r + [" "] * (n_cols - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "|" + "|".join(["---"] * n_cols) + "|"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)
